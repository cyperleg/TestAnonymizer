import os
import logging
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl

# Configure module-level logging
logger = logging.getLogger("app")

class FileTextExtractor:
    """
    Class for extracting text from common file types: txt, pdf, docx, pptx, xlsx.
    """

    @classmethod
    def extract(cls, file_path: str, encoding: str = 'utf-8') -> str:
        """
        Determine file format by extension and delegate to the corresponding extraction method.

        Args:
            file_path: Path to the file to extract text from.
            encoding: Encoding for text files (default: 'utf-8').

        Returns:
            All extracted text as a single string.

        Raises:
            ValueError: If the file extension is not supported.
        """
        logger.info(f"Starting text extraction from: {file_path}")
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.txt':
            logger.debug("Detected TXT file format.")
            text = cls._extract_txt(file_path, encoding)
        elif ext == '.pdf':
            logger.debug("Detected PDF file format.")
            text = cls._extract_pdf(file_path)
        elif ext == '.docx':
            logger.debug("Detected DOCX file format.")
            text = cls._extract_docx(file_path)
        elif ext == '.pptx':
            logger.debug("Detected PPTX file format.")
            text = cls._extract_pptx(file_path)
        elif ext == '.xlsx':
            logger.debug("Detected XLSX file format.")
            text = cls._extract_xlsx(file_path)
        else:
            logger.error(f"Unsupported file extension: {ext}")
            raise ValueError(f"Unsupported file extension: {ext}")

        logger.info(f"Completed extraction ({len(text)} characters)")
        return text

    @classmethod
    def _extract_txt(cls, path: str, encoding: str) -> str:
        """
        Read and return the content of a plain text file.
        """
        logger.debug(f"Reading TXT file: {path}")
        with open(path, 'r', encoding=encoding, errors='ignore') as f:
            return f.read()

    @classmethod
    def _extract_pdf(cls, path: str) -> str:
        """
        Extract text from each page of a PDF file.
        """
        logger.debug(f"Reading PDF file: {path}")
        text_pages = []
        with open(path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page_number, page in enumerate(reader.pages, start=1):
                page_text = page.extract_text() or ""
                logger.debug(f"Extracted page {page_number}, {len(page_text)} chars")
                text_pages.append(page_text)
        return "\n".join(text_pages)

    @classmethod
    def _extract_docx(cls, path: str) -> str:
        """
        Extract text from all paragraphs in a DOCX document.
        """
        logger.debug(f"Reading DOCX file: {path}")
        doc = Document(path)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        logger.debug(f"Extracted {len(paragraphs)} paragraphs from DOCX")
        return "\n".join(paragraphs)

    @classmethod
    def _extract_pptx(cls, path: str) -> str:
        """
        Extract text from all text-containing shapes in a PPTX presentation.
        """
        logger.debug(f"Reading PPTX file: {path}")
        prs = Presentation(path)
        slides_text = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slides_text.append(shape.text)
            logger.debug(f"Extracted slide {slide_index}, shapes: {len(slides_text)} total pieces so far")
        return "\n".join(slides_text)

    @classmethod
    def _extract_xlsx(cls, path: str) -> str:
        """
        Extract text from all cells in an XLSX workbook.
        """
        logger.debug(f"Reading XLSX file: {path}")
        wb = openpyxl.load_workbook(path, data_only=True)
        sheets_text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) for cell in row if cell is not None)
                if row_text:
                    sheets_text.append(row_text)
            logger.debug(f"Extracted sheet '{sheet.title}'")
        return "\n".join(sheets_text)

if __name__ == "__main__":
    extractor = FileTextExtractor()
    for ext in [
        "txt",
        "pdf",
        "docx",
        "pptx",
        "xlsx"
    ]:
        try:
            content = extractor.extract("files/sample." + ext)
            print(f"--- sample.{ext} ---\n{content[:200]}...\n")
        except Exception as err:
            print(f"Failed to process sample.{ext}: {err}")
